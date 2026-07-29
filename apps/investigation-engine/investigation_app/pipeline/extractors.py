"""Real, offline text + metadata extraction for evidence files.

Every extractor here is:

* **Offline** - no network, no cloud services. OCR prefers the portal's
  shared RapidOCR/ONNX engine and falls back to local Tesseract when present;
  document parsing uses local Python libraries only.
* **Optional / lazily loaded** - heavy libraries (pypdf, python-docx,
  openpyxl, python-pptx, Pillow, rapidocr, onnxruntime, pytesseract) are imported *inside* the
  functions. If a library (or OCR engine) is not installed, the
  extractor returns ``""`` / ``{}`` instead of raising, so the pipeline keeps
  running and simply gets less text for that file. This preserves the
  project's "degrade gracefully, never fail the pipeline" contract.
* **Bounded** - output text is capped so a huge document cannot exhaust RAM.

The module exposes two entry points used by :mod:`investigation_app.pipeline.metadata`:
``extract_text(path, kind, mime)`` and ``extract_metadata(path, kind)``.
"""
from __future__ import annotations

import base64
import io
import json
import logging
import os
import shutil
import tarfile
import urllib.error
import urllib.request
import zipfile
from typing import Any, Callable, Dict, Optional

_logger = logging.getLogger("iie.extractors")

# Bound extracted text so one file can never exhaust memory / the FTS index.
MAX_TEXT_CHARS = int(os.environ.get("IIE_MAX_TEXT_CHARS", "64000000"))
# OCR is the slowest stage; cap the number of PDF pages/image passes.
MAX_PDF_PAGES = 200
MAX_OCR_PAGES = 30


def _truncate(text: str) -> str:
    if text and len(text) > MAX_TEXT_CHARS:
        return text[:MAX_TEXT_CHARS]
    return text or ""


def _decode_text_bytes(data: bytes) -> str:
    """Decode evidence bytes without executing/interpreting the file."""
    if not data:
        return ""
    for enc in ("utf-8-sig", "utf-16", "cp1252", "latin-1"):
        try:
            return data.decode(enc, errors="replace")
        except Exception:
            continue
    return data.decode("utf-8", errors="replace")


def _bytes_look_text(data: bytes) -> bool:
    if not data:
        return True
    if b"\x00" in data:
        return False
    printable = sum(1 for b in data[:8192] if b in b"\t\n\r\f\b" or 32 <= b <= 126 or b >= 128)
    return printable / max(min(len(data), 8192), 1) >= 0.85



def _effective_ext(path: str, filename_hint: str | None = None) -> str:
    """Return extension from original filename when stored path is hash-only."""
    return (os.path.splitext(filename_hint or "")[1] or os.path.splitext(path)[1]).lower()


def _office_ext_from_content(path: str, filename_hint: str | None = None) -> str:
    """Detect Office subtype even when stored evidence has no extension."""
    ext = _effective_ext(path, filename_hint)
    if ext in {".docx", ".xlsx", ".pptx", ".xls"}:
        return ext
    try:
        if zipfile.is_zipfile(path):
            with zipfile.ZipFile(path) as zf:
                names = set(zf.namelist())
                if any(n.startswith("xl/") for n in names):
                    return ".xlsx"
                if any(n.startswith("word/") for n in names):
                    return ".docx"
                if any(n.startswith("ppt/") for n in names):
                    return ".pptx"
    except Exception:
        pass
    return ext

# --- Plain text ----------------------------------------------------------

def _read_plain(path: str) -> str:
    try:
        with open(path, "rb") as fh:
            return _decode_text_bytes(fh.read(MAX_TEXT_CHARS))
    except OSError:
        return ""


# --- HTML ----------------------------------------------------------------

def _strip_html(path: str) -> str:
    """Extract visible text from HTML using only the stdlib html parser."""
    from html.parser import HTMLParser

    class _Text(HTMLParser):
        def __init__(self) -> None:
            super().__init__()
            self.parts: list[str] = []
            self._skip = 0

        def handle_starttag(self, tag, attrs):
            if tag in {"script", "style"}:
                self._skip += 1

        def handle_endtag(self, tag):
            if tag in {"script", "style"} and self._skip:
                self._skip -= 1

        def handle_data(self, data):
            if not self._skip and data.strip():
                self.parts.append(data)

    raw = _read_plain(path)
    if not raw:
        return ""
    try:
        parser = _Text()
        parser.feed(raw)
        return " ".join(parser.parts)
    except Exception:  # noqa: BLE001 - malformed HTML must not break the run
        return raw


# --- Email (.eml) --------------------------------------------------------

def _extract_eml(path: str) -> str:
    """Extract headers + text body from an RFC822 email (stdlib only)."""
    try:
        from email import policy
        from email.parser import BytesParser

        with open(path, "rb") as fh:
            msg = BytesParser(policy=policy.default).parse(fh)
        header_keys = ("From", "To", "Cc", "Subject", "Date")
        lines = [f"{k}: {msg[k]}" for k in header_keys if msg[k]]
        body = ""
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == "text/plain":
                    body += part.get_content()
        else:
            if msg.get_content_type() == "text/plain":
                body = msg.get_content()
        return _truncate("\n".join(lines) + "\n\n" + (body or ""))
    except Exception:  # noqa: BLE001
        _logger.info("eml extract failed for %s", path)
        return ""


# --- PDF -----------------------------------------------------------------

def _extract_pdf(path: str) -> str:
    """Extract embedded text from a PDF; OCR pages that have no text layer."""
    text_parts: list[str] = []
    try:
        from pypdf import PdfReader  # optional dep

        reader = PdfReader(path)
        pages = reader.pages[:MAX_PDF_PAGES]
        empty_pages = []
        for i, page in enumerate(pages):
            try:
                t = page.extract_text() or ""
            except Exception:  # noqa: BLE001 - one bad page must not abort
                t = ""
            if t.strip():
                text_parts.append(t)
            else:
                empty_pages.append(i)
        # Scanned PDFs have no text layer -> OCR the image-only pages.
        if empty_pages:
            ocr = _ocr_pdf_pages(path, empty_pages[:MAX_OCR_PAGES])
            if ocr:
                text_parts.append(ocr)
    except Exception:  # noqa: BLE001 - pypdf missing or file unreadable
        _logger.info("pdf text extract unavailable for %s", path)
    return _truncate("\n".join(text_parts))


def _ocr_pdf_pages(path: str, page_indexes: list[int]) -> str:
    """OCR specific PDF pages by rasterising with PyMuPDF.

    Prefer the shared RapidOCR engine already used by the AI Assistant; fall
    back to Tesseract only when RapidOCR is unavailable or blank.
    """
    try:
        import fitz  # PyMuPDF, optional
        from PIL import Image  # optional
        import io

        doc = fitz.open(path)
        out: list[str] = []
        try:
            for idx in page_indexes:
                if idx >= doc.page_count:
                    continue
                pix = doc.load_page(idx).get_pixmap(dpi=220)
                img = Image.open(io.BytesIO(pix.tobytes("png")))
                text, diag = _run_shared_rapidocr(img)
                if not text.strip():
                    text, diag = _run_tesseract_ocr(img)
                if text.strip():
                    out.append(f"[PDF OCR page {idx + 1}]\n{text.strip()}")
                else:
                    out.append(f"[PDF OCR page {idx + 1} diagnostics] {diag}")
        finally:
            doc.close()
        return "\n".join(p for p in out if p.strip())
    except Exception as exc:  # noqa: BLE001 - any missing OCR dep -> skip silently
        _logger.info("pdf OCR unavailable for %s: %s", path, exc)
        return ""


# --- Images (OCR-first, optional explicit local vision override) ----------


def _run_shared_rapidocr(img) -> tuple[str, str]:
    """Run the same RapidOCR/ONNX-style OCR stack used by the AI Assistant.

    The AI Assistant's original module lives under an ``app`` package, which
    collides with the Investigation Engine's own ``app`` package at runtime.
    Therefore the reusable engine is exposed through the neutral ``shared`` package.
    """
    try:
        from shared.offline_ocr import rapidocr_pil_image
    except Exception as exc:
        return "", f"shared RapidOCR helper unavailable: {exc}"
    return rapidocr_pil_image(img)

def _find_tesseract_cmd() -> str:
    """Find a local Tesseract binary without requiring users to edit PATH."""
    env = os.environ.get("TESSERACT_CMD") or os.environ.get("TESSERACT_PATH")
    candidates = [
        env,
        shutil.which("tesseract"),
        r"C:\Program Files\Tesseract-OCR\tesseract.exe",
        r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
        "/usr/bin/tesseract",
        "/usr/local/bin/tesseract",
    ]
    for item in candidates:
        if item and os.path.isfile(item):
            return item
    return ""


def _prepare_ocr_variants(img):
    """Return screenshot-friendly image variants for Tesseract OCR."""
    try:
        from PIL import ImageFilter, ImageOps
    except Exception:
        return [img]
    variants = []
    base = img.convert("RGB")
    # Upscale small phone screenshots; Tesseract performs much better around
    # 1800-2400px width for chat/profile UI text.
    try:
        w, h = base.size
        if w < 1800:
            scale = min(3.0, 1800 / max(w, 1))
            base = base.resize((int(w * scale), int(h * scale)))
    except Exception:
        pass
    variants.append(base)
    try:
        gray = ImageOps.grayscale(base)
        gray = ImageOps.autocontrast(gray)
        variants.append(gray)
        variants.append(gray.filter(ImageFilter.SHARPEN))
        # High-contrast threshold variant helps WhatsApp/Instagram screenshots.
        variants.append(gray.point(lambda p: 255 if p > 170 else 0))
    except Exception:
        pass
    return variants


def _run_tesseract_ocr(img) -> tuple[str, str]:
    """OCR an image using local Tesseract. Returns (text, diagnostic)."""
    try:
        import pytesseract  # optional
    except Exception as exc:
        return "", f"pytesseract Python package is not installed: {exc}"

    cmd = _find_tesseract_cmd()
    if cmd:
        try:
            pytesseract.pytesseract.tesseract_cmd = cmd
        except Exception:
            pass

    texts: list[str] = []
    errors: list[str] = []
    configs = ("--oem 3 --psm 6", "--oem 3 --psm 11")
    # Try English first. If Hindi/Marathi traineddata exists, users can set
    # IIE_OCR_LANG=eng+hin+mar. Bad language config is caught and retried.
    langs = [os.environ.get("IIE_OCR_LANG", "eng"), "eng"]
    tried = set()
    for variant in _prepare_ocr_variants(img):
        for lang in langs:
            for cfg in configs:
                key = (lang, cfg, id(variant))
                if key in tried:
                    continue
                tried.add(key)
                try:
                    text = pytesseract.image_to_string(variant, lang=lang, config=cfg) or ""
                    if text.strip():
                        texts.append(text.strip())
                except Exception as exc:  # noqa: BLE001
                    msg = str(exc).splitlines()[0][:180]
                    if msg not in errors:
                        errors.append(msg)
    if texts:
        # Use the richest output but keep all distinct lines from alternatives.
        lines: list[str] = []
        seen = set()
        for text in sorted(texts, key=len, reverse=True)[:3]:
            for line in text.splitlines():
                clean = line.strip()
                key = clean.lower()
                if clean and key not in seen:
                    seen.add(key)
                    lines.append(clean)
        return "\n".join(lines), "ok"
    return "", "; ".join(errors[:3]) or "Tesseract returned no readable text"


def _ollama_vision_extract(path: str) -> tuple[str, str]:
    """Optional local-only vision fallback using an installed Ollama VLM."""
    if os.environ.get("IIE_ENABLE_VISION_FALLBACK", "0").lower() in {"0", "false", "no", "off", ""}:
        return "", "disabled by IIE_ENABLE_VISION_FALLBACK"
    model = os.environ.get("IIE_VISION_MODEL", "").strip()
    if not model:
        return "", "disabled: no IIE_VISION_MODEL configured"
    base = os.environ.get("OLLAMA_HOST") or os.environ.get("IIE_OLLAMA_URL") or "http://127.0.0.1:11434"
    base = base.rstrip("/")
    try:
        # Probe tags first so missing Ollama/model fails quickly instead of
        # hanging the evidence worker.
        with urllib.request.urlopen(f"{base}/api/tags", timeout=2) as resp:
            tags = json.loads(resp.read().decode("utf-8", errors="replace"))
        names = {m.get("name") for m in tags.get("models", []) if isinstance(m, dict)}
        if model not in names and not any(str(n or "").startswith(model.split(":")[0] + ":") for n in names):
            return "", f"Ollama vision model '{model}' is not installed"
        with open(path, "rb") as fh:
            img64 = base64.b64encode(fh.read()).decode("ascii")
        prompt = (
            "You are an offline digital evidence evidence extractor. Read this screenshot/image. "
            "Return only facts visible in the image. Extract visible text, platform/app, "
            "usernames/handles, names, phone numbers, emails, URLs/domains, UPI IDs, bank "
            "accounts, IFSC, UTR/reference IDs, amounts, dates/times, messages, QR/payment "
            "payloads if visible, and suspicious/risk indicators. Do not guess. If unreadable, say unreadable."
        )
        payload = json.dumps({
            "model": model,
            "prompt": prompt,
            "images": [img64],
            "stream": False,
            "options": {"temperature": 0.0},
        }).encode("utf-8")
        req = urllib.request.Request(f"{base}/api/generate", data=payload, headers={"Content-Type": "application/json"}, method="POST")
        timeout = int(os.environ.get("IIE_VISION_TIMEOUT", "180"))
        with urllib.request.urlopen(req, timeout=timeout) as resp:
            data = json.loads(resp.read().decode("utf-8", errors="replace"))
        text = str(data.get("response") or "").strip()
        if text:
            return text, "ok"
        return "", "Ollama vision returned no text"
    except (urllib.error.URLError, urllib.error.HTTPError, OSError, ValueError) as exc:
        return "", f"Ollama vision unavailable: {exc}"


def _extract_image(path: str) -> str:
    """OCR/image intelligence for screenshots with explicit diagnostics."""
    parts: list[str] = []
    diagnostics: list[str] = []
    try:
        from PIL import Image  # optional
    except Exception as exc:
        return _truncate(f"[Image evidence]\nImage decode unavailable: Pillow is not installed ({exc}).")

    try:
        with Image.open(path) as img:
            try:
                parts.append(f"[Image evidence] format={img.format or 'unknown'} size={img.size[0]}x{img.size[1]} mode={img.mode}")
            except Exception:
                parts.append("[Image evidence]")

            # First use the same RapidOCR/ONNX OCR stack already present in
            # the AI Assistant. This avoids requiring the external Tesseract
            # binary for normal screenshot evidence.
            ocr_text, ocr_diag = _run_shared_rapidocr(img)
            ocr_source = "RapidOCR"
            if not ocr_text.strip():
                diagnostics.append("RapidOCR: " + ocr_diag)
                ocr_text, ocr_diag = _run_tesseract_ocr(img)
                ocr_source = "Tesseract"

            if ocr_text.strip():
                parts.append(f"[OCR text via {ocr_source}]\n" + ocr_text.strip())
            else:
                diagnostics.append("Tesseract: " + ocr_diag)

            # Optional QR decoding. If pyzbar/libzbar is absent, skip without
            # affecting OCR/vision or the rest of the pipeline.
            try:
                from pyzbar.pyzbar import decode  # optional
                for idx, item in enumerate(decode(img), start=1):
                    try:
                        payload = item.data.decode("utf-8", errors="replace")
                    except Exception:
                        payload = str(item.data)
                    if payload:
                        parts.append(f"QR payload {idx}: {payload}")
            except Exception as exc:
                diagnostics.append("QR decode: optional pyzbar/libzbar unavailable")
    except Exception as exc:  # noqa: BLE001 - bad/corrupt image
        return _truncate(f"[Image evidence]\nImage file could not be opened: {exc}")

    # OCR is the default image intelligence path. Direct vision analysis is
    # disabled unless explicitly enabled via IIE_ENABLE_VISION_FALLBACK and
    # IIE_IMAGE_VISION_MODE. This keeps qwen3:8b text-only deployments stable.
    ocr_chars = len(ocr_text.strip()) if "ocr_text" in locals() and isinstance(ocr_text, str) else 0
    vision_mode = os.environ.get("IIE_IMAGE_VISION_MODE", "disabled").strip().lower()
    min_chars = int(os.environ.get("IIE_VISION_MIN_OCR_CHARS", "180"))
    has_ocr_text = any("[OCR text via" in p for p in parts)
    needs_vision = (
        vision_mode in {"1", "true", "yes", "always"}
        or (vision_mode == "auto" and ocr_chars < min_chars)
    )
    if needs_vision:
        vision_text, vision_diag = _ollama_vision_extract(path)
        if vision_text.strip():
            parts.append("[Local vision analysis]\n" + vision_text.strip())
        else:
            diagnostics.append("Vision: " + vision_diag)

    if diagnostics:
        parts.append("[Image processing diagnostics]\n" + "\n".join(f"- {d}" for d in diagnostics[:6]))
    return _truncate("\n".join(parts))


# --- Office documents ----------------------------------------------------

def _extract_office(path: str, filename_hint: str | None = None, progress_callback: Optional[Callable[[float, str, int, int], None]] = None) -> str:
    ext = _office_ext_from_content(path, filename_hint)
    try:
        if ext == ".docx":
            import docx  # python-docx, optional

            document = docx.Document(path)
            parts = [p.text for p in document.paragraphs if p.text]
            for table in document.tables:
                for row in table.rows:
                    parts.append(" ".join(c.text for c in row.cells))
            return _truncate("\n".join(parts))
        if ext == ".xlsx":
            import openpyxl  # optional

            # Stored evidence files are named by SHA-256 without extension;
            # openpyxl rejects extensionless paths, but accepts a binary file
            # object. Keep the handle open until the workbook is closed.
            fh = open(path, "rb")
            try:
                wb = openpyxl.load_workbook(fh, read_only=True, data_only=True)
                parts: list[str] = []
                try:
                    max_rows = int(os.environ.get("IIE_MAX_XLSX_ROWS", "1000000"))
                    approx_chars = 0
                    emitted_rows = 0
                    scanned_rows = 0
                    stopped = False
                    try:
                        total_rows = sum(int(getattr(ws, "max_row", 0) or 0) for ws in wb.worksheets)
                    except Exception:
                        total_rows = max_rows
                    total_rows = max(1, min(total_rows or max_rows, max_rows))
                    progress_step = max(25, total_rows // 100)
                    if progress_callback:
                        progress_callback(0.0, "Opening Excel workbook", 0, total_rows)
                    for ws in wb.worksheets:
                        if stopped:
                            break
                        line = f"[Sheet: {ws.title}]"
                        parts.append(line); approx_chars += len(line) + 1
                        rows = ws.iter_rows(values_only=True)
                        headers = []
                        for row_idx, row in enumerate(rows, start=1):
                            scanned_rows += 1
                            if progress_callback and (scanned_rows == 1 or scanned_rows % progress_step == 0):
                                pct = min(99.0, (scanned_rows / max(total_rows, 1)) * 100.0)
                                progress_callback(pct, f"Reading Excel rows {min(scanned_rows, total_rows):,}/{total_rows:,}", min(scanned_rows, total_rows), total_rows)
                            values = ["" if c is None else str(c).strip() for c in row]
                            if not any(values):
                                continue
                            # First non-empty row becomes header context for structured
                            # transaction workbooks. We still include raw cells for odd files.
                            if not headers:
                                headers = [v or f"Column {i+1}" for i, v in enumerate(values)]
                                line = f"Sheet {ws.title} header | " + " | ".join(headers)
                                parts.append(line); approx_chars += len(line) + 1
                                continue
                            pairs = []
                            for i, val in enumerate(values):
                                if not val:
                                    continue
                                header = headers[i] if i < len(headers) and headers[i] else f"Column {i+1}"
                                pairs.append(f"{header}: {val}")
                            if pairs:
                                line = f"Sheet {ws.title} row {row_idx} | " + " | ".join(pairs)
                                parts.append(line); approx_chars += len(line) + 1; emitted_rows += 1
                            if emitted_rows >= max_rows or approx_chars >= MAX_TEXT_CHARS:
                                parts.append(f"[XLSX extraction capped after {emitted_rows} data rows / {approx_chars} chars. Increase IIE_MAX_XLSX_ROWS only if needed.]")
                                if progress_callback:
                                    progress_callback(100.0, f"Excel extraction capped after {emitted_rows:,} data rows", emitted_rows, max(emitted_rows, total_rows))
                                stopped = True
                                break
                finally:
                    wb.close()
                if progress_callback:
                    progress_callback(100.0, f"Excel extraction read {emitted_rows:,} data rows", emitted_rows, max(emitted_rows, total_rows))
                return _truncate("\n".join(parts))
            finally:
                fh.close()
        if ext == ".xls":
            # Optional fallback for legacy Excel. No hard dependency: if pandas/xlrd
            # is absent, the pipeline still continues with metadata-only evidence.
            try:
                import pandas as pd  # optional
                xls = pd.ExcelFile(path)
                parts = []
                for sheet in xls.sheet_names:
                    df = xls.parse(sheet, dtype=str).fillna("")
                    parts.append(f"[Sheet: {sheet}]")
                    for idx, row in df.iterrows():
                        pairs = [f"{col}: {str(row[col]).strip()}" for col in df.columns if str(row[col]).strip()]
                        if pairs:
                            parts.append(f"Sheet {sheet} row {int(idx)+2} | " + " | ".join(pairs))
                return _truncate("\n".join(parts))
            except Exception:
                return ""
        if ext == ".pptx":
            from pptx import Presentation  # optional

            prs = Presentation(path)
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        parts.append(shape.text_frame.text)
            return _truncate("\n".join(p for p in parts if p))
    except Exception:  # noqa: BLE001 - missing lib / legacy .doc / bad file
        _logger.info("office extract unavailable for %s", path)
    return ""


# --- Archives ------------------------------------------------------------

def _extract_archive(path: str, filename_hint: str | None = None) -> str:
    """Index archive manifest and small text members without extracting to disk.

    This is safe evidence triage, not execution. Binary members are listed by
    name/size only. Text-like members get a bounded snippet so entities inside
    logs/scripts/configs packaged in ZIP/TAR evidence are still searchable.
    """
    ext = _effective_ext(path, filename_hint)
    parts: list[str] = []
    max_members = 250
    max_member_bytes = 80_000
    try:
        if zipfile.is_zipfile(path):
            with zipfile.ZipFile(path) as zf:
                for idx, info in enumerate(zf.infolist()[:max_members], start=1):
                    if info.is_dir():
                        continue
                    parts.append(f"Archive member {idx}: {info.filename} size {info.file_size}")
                    if info.file_size <= max_member_bytes:
                        data = zf.read(info.filename)
                        if _bytes_look_text(data):
                            parts.append(f"--- text from {info.filename} ---\n" + _decode_text_bytes(data[:max_member_bytes]))
                return _truncate("\n".join(parts))
        if tarfile.is_tarfile(path):
            with tarfile.open(path) as tf:
                for idx, member in enumerate([m for m in tf.getmembers() if m.isfile()][:max_members], start=1):
                    parts.append(f"Archive member {idx}: {member.name} size {member.size}")
                    if member.size <= max_member_bytes:
                        fh = tf.extractfile(member)
                        if fh:
                            data = fh.read(max_member_bytes)
                            if _bytes_look_text(data):
                                parts.append(f"--- text from {member.name} ---\n" + _decode_text_bytes(data))
                return _truncate("\n".join(parts))
    except Exception:  # noqa: BLE001
        _logger.info("archive extract unavailable for %s", path)
    return ""


# --- Metadata ------------------------------------------------------------

def _image_metadata(path: str) -> Dict[str, Any]:
    """Image metadata incl. dimensions plus EXIF/GPS when present."""
    meta: Dict[str, Any] = {}
    try:
        from PIL import Image, ExifTags  # optional

        with Image.open(path) as img:
            try:
                meta["format"] = str(img.format or "")
                meta["width"] = int(img.size[0])
                meta["height"] = int(img.size[1])
                meta["mode"] = str(img.mode or "")
            except Exception:
                pass
            exif = getattr(img, "_getexif", lambda: None)()
            if not exif:
                return meta
            tags = {ExifTags.TAGS.get(k, k): v for k, v in exif.items()}
            for key in ("DateTimeOriginal", "DateTime", "Make", "Model", "Artist"):
                if tags.get(key):
                    meta[key] = str(tags[key])
            gps = tags.get("GPSInfo")
            if gps:
                coords = _decode_gps(gps, ExifTags)
                if coords:
                    meta["gps"] = coords
    except Exception:  # noqa: BLE001
        pass
    return meta


def _decode_gps(gps: Dict[Any, Any], ExifTags) -> Dict[str, float]:
    try:
        g = {ExifTags.GPSTAGS.get(k, k): v for k, v in gps.items()}

        def _to_deg(value) -> float:
            d, m, s = value
            return float(d) + float(m) / 60.0 + float(s) / 3600.0

        lat = _to_deg(g["GPSLatitude"])
        lon = _to_deg(g["GPSLongitude"])
        if g.get("GPSLatitudeRef") == "S":
            lat = -lat
        if g.get("GPSLongitudeRef") == "W":
            lon = -lon
        return {"lat": round(lat, 6), "lon": round(lon, 6)}
    except Exception:  # noqa: BLE001
        return {}


def _pdf_metadata(path: str) -> Dict[str, Any]:
    meta: Dict[str, Any] = {}
    try:
        from pypdf import PdfReader  # optional

        reader = PdfReader(path)
        info = reader.metadata or {}
        for key in ("/Author", "/Title", "/Producer", "/CreationDate", "/ModDate"):
            if info.get(key):
                meta[key.strip("/").lower()] = str(info.get(key))
        meta["pages"] = len(reader.pages)
    except Exception:  # noqa: BLE001
        pass
    return meta


def _office_metadata(path: str, filename_hint: str | None = None) -> Dict[str, Any]:
    meta: Dict[str, Any] = {}
    ext = _office_ext_from_content(path, filename_hint)
    try:
        cp = None
        if ext == ".docx":
            import docx
            cp = docx.Document(path).core_properties
        elif ext == ".xlsx":
            import openpyxl
            fh = open(path, "rb")
            try:
                wb = openpyxl.load_workbook(fh, read_only=True)
                cp = wb.properties
            finally:
                try:
                    wb.close()
                except Exception:
                    pass
                fh.close()
        elif ext == ".pptx":
            from pptx import Presentation
            cp = Presentation(path).core_properties
        if cp is not None:
            for attr in ("author", "last_modified_by", "created", "modified", "title"):
                val = getattr(cp, attr, None)
                if val:
                    meta[attr] = str(val)
    except Exception:  # noqa: BLE001
        pass
    return meta


# --- Public API ----------------------------------------------------------

def extract_text(path: str, kind: str, mime: str = "", filename_hint: str | None = None, progress_callback: Optional[Callable[[float, str, int, int], None]] = None) -> str:
    """Best-effort plain text for any supported evidence kind."""
    ext = _effective_ext(path, filename_hint)
    try:
        if ext in {".html", ".htm", ".xhtml", ".svg"}:
            return _truncate(_strip_html(path))
        if ext == ".eml":
            return _extract_eml(path)
        if kind == "text":
            return _truncate(_read_plain(path))
        if kind == "pdf":
            return _extract_pdf(path)
        if kind == "image":
            return _extract_image(path)
        if kind == "office":
            return _extract_office(path, filename_hint=filename_hint, progress_callback=progress_callback)
        if kind == "archive":
            return _extract_archive(path, filename_hint=filename_hint)
        if kind == "binary":
            # Last-resort text sniff: extensionless scripts or tool dumps can be
            # scanned, but real binaries stay metadata/hash-only.
            try:
                with open(path, "rb") as fh:
                    data = fh.read(min(os.path.getsize(path), MAX_TEXT_CHARS))
                if _bytes_look_text(data):
                    return _truncate(_decode_text_bytes(data))
            except OSError:
                return ""
    except Exception:  # noqa: BLE001 - top-level guard; never fail the caller
        _logger.exception("extract_text failed for %s", path)
    return ""


def extract_metadata(path: str, kind: str, filename_hint: str | None = None) -> Dict[str, Any]:
    """Best-effort rich metadata (EXIF/GPS, PDF/Office properties)."""
    try:
        if kind == "image":
            return _image_metadata(path)
        if kind == "pdf":
            return _pdf_metadata(path)
        if kind == "office":
            return _office_metadata(path, filename_hint=filename_hint)
    except Exception:  # noqa: BLE001
        _logger.exception("extract_metadata failed for %s", path)
    return {}
